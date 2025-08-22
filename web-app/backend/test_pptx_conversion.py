#!/usr/bin/env python3
"""Test actual PPTX conversion with markitdown"""

from pptx import Presentation
from markitdown import MarkItDown
import tempfile
import os

def test_pptx_conversion():
    """Test PPTX file conversion"""
    try:
        # Create a simple PPTX file using python-pptx
        prs = Presentation()
        
        # Add a title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Test Presentation"
        subtitle.text = "This is a test subtitle"
        
        # Add a content slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Test Slide"
        
        tf = body_shape.text_frame
        tf.text = "First bullet point"
        
        p = tf.add_paragraph()
        p.text = "Second bullet point"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = "Third bullet point"
        p.level = 0
        
        # Save the presentation to a temporary file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
            pptx_path = tmp_file.name
            prs.save(pptx_path)
            print(f"✓ Created test PPTX file: {pptx_path}")
        
        # Now test conversion with markitdown
        md = MarkItDown()
        result = md.convert(pptx_path)
        
        print("✓ Conversion successful!")
        print("\nConverted content:")
        print("-" * 40)
        print(result.text_content[:500])  # Print first 500 characters
        print("-" * 40)
        
        # Clean up
        os.unlink(pptx_path)
        print("✓ Cleaned up test file")
        
        return True
        
    except Exception as e:
        print(f"✗ Error during conversion: {e}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    success = test_pptx_conversion()
    exit(0 if success else 1)