#!/usr/bin/env python3
"""Test FastAPI server with PPTX conversion"""

import asyncio
import httpx
import tempfile
import os
from pptx import Presentation

async def test_api_conversion():
    """Test conversion through API"""
    
    # Create a test PPTX file
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "API Test Presentation"
    subtitle.text = "Testing via FastAPI"
    
    # Add content slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Content Slide"
    tf = body_shape.text_frame
    tf.text = "First point"
    p = tf.add_paragraph()
    p.text = "Second point"
    
    # Save to temp file
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
        pptx_path = tmp_file.name
        prs.save(pptx_path)
        print(f"Created test file: {pptx_path}")
    
    try:
        async with httpx.AsyncClient() as client:
            # Test file upload and conversion
            with open(pptx_path, 'rb') as f:
                files = {"file": ("test.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
                response = await client.post(
                    "http://localhost:8000/api/v1/conversion/upload",
                    files=files
                )
            
            print(f"\nAPI Response Status: {response.status_code}")
            
            if response.status_code == 200:
                result = response.json()
                print(f"Conversion Result:")
                print(f"  Status: {result.get('status')}")
                print(f"  Input: {result.get('input_file')}")
                print(f"  Output: {result.get('output_file')}")
                
                if result.get('markdown_content'):
                    print(f"\nMarkdown Content (first 500 chars):")
                    print("-" * 40)
                    print(result['markdown_content'][:500])
                    print("-" * 40)
                    
                return True
            else:
                print(f"Error: {response.text}")
                return False
                
    except Exception as e:
        print(f"Error: {e}")
        return False
    finally:
        # Clean up
        os.unlink(pptx_path)
        print("\nCleaned up test file")

if __name__ == "__main__":
    print("Testing API server (make sure it's running on port 8000)")
    success = asyncio.run(test_api_conversion())
    exit(0 if success else 1)