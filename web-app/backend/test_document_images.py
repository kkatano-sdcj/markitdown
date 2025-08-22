#!/usr/bin/env python3
"""
Test script to create sample documents with embedded images
and test the image extraction functionality
"""
import os
import sys
from PIL import Image, ImageDraw, ImageFont
import io

# Add app to path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

def create_test_image_with_text(text: str, size=(400, 200)):
    """Create a test image with text"""
    # Create white background
    img = Image.new('RGB', size, color='white')
    draw = ImageDraw.Draw(img)
    
    # Draw text (use default font)
    draw.text((50, 50), text, fill='black')
    
    # Draw a border
    draw.rectangle([0, 0, size[0]-1, size[1]-1], outline='blue', width=2)
    
    return img

def create_docx_with_images():
    """Create a Word document with embedded images"""
    try:
        from docx import Document
        from docx.shared import Inches
        
        doc = Document()
        doc.add_heading('Test Document with Images', 0)
        
        doc.add_paragraph('This document contains embedded images for testing.')
        
        # Create and add first image
        img1 = create_test_image_with_text("First Image: Hello World!")
        img1_bytes = io.BytesIO()
        img1.save(img1_bytes, format='PNG')
        img1_bytes.seek(0)
        
        doc.add_heading('Section 1: First Image', level=1)
        doc.add_paragraph('Below is an embedded image with text:')
        doc.add_picture(img1_bytes, width=Inches(3))
        
        # Create and add second image
        img2 = create_test_image_with_text("Second Image: Test OCR")
        img2_bytes = io.BytesIO()
        img2.save(img2_bytes, format='PNG')
        img2_bytes.seek(0)
        
        doc.add_heading('Section 2: Second Image', level=1)
        doc.add_paragraph('Another embedded image:')
        doc.add_picture(img2_bytes, width=Inches(3))
        
        # Save document
        doc.save('test_document_with_images.docx')
        print("Created: test_document_with_images.docx")
        return True
        
    except ImportError:
        print("python-docx not installed")
        return False

def create_pptx_with_images():
    """Create a PowerPoint presentation with embedded images"""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation()
        
        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Test Presentation"
        subtitle.text = "With Embedded Images"
        
        # Slide with image
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
        slide.shapes.title.text = "Slide with Image"
        
        # Create and add image
        img = create_test_image_with_text("PowerPoint Image: OCR Test")
        img_bytes = io.BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        
        left = Inches(2)
        top = Inches(2)
        slide.shapes.add_picture(img_bytes, left, top, width=Inches(4))
        
        # Another slide with different image
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Another Image"
        
        img2 = create_test_image_with_text("Slide 2: Text Recognition")
        img2_bytes = io.BytesIO()
        img2.save(img2_bytes, format='PNG')
        img2_bytes.seek(0)
        
        slide.shapes.add_picture(img2_bytes, left, top, width=Inches(4))
        
        # Save presentation
        prs.save('test_presentation_with_images.pptx')
        print("Created: test_presentation_with_images.pptx")
        return True
        
    except ImportError:
        print("python-pptx not installed")
        return False

def test_extraction():
    """Test the extraction functionality"""
    from app.services.enhanced_conversion_service import EnhancedConversionService
    import asyncio
    
    service = EnhancedConversionService()
    
    # Test with Word document if it exists
    if os.path.exists('test_document_with_images.docx'):
        print("\nTesting Word document extraction...")
        result = asyncio.run(service.convert_file_enhanced(
            'test_document_with_images.docx',
            'test_docx_output.md'
        ))
        if result.status.value == 'completed':
            print("✓ Word document processed successfully")
            print(f"  Output saved to: {result.output_file}")
        else:
            print(f"✗ Failed: {result.error_message}")
    
    # Test with PowerPoint if it exists
    if os.path.exists('test_presentation_with_images.pptx'):
        print("\nTesting PowerPoint extraction...")
        result = asyncio.run(service.convert_file_enhanced(
            'test_presentation_with_images.pptx',
            'test_pptx_output.md'
        ))
        if result.status.value == 'completed':
            print("✓ PowerPoint processed successfully")
            print(f"  Output saved to: {result.output_file}")
        else:
            print(f"✗ Failed: {result.error_message}")

if __name__ == "__main__":
    print("Creating test documents with embedded images...\n")
    
    # Create test documents
    docx_created = create_docx_with_images()
    pptx_created = create_pptx_with_images()
    
    if docx_created or pptx_created:
        print("\nTesting image extraction and OCR...")
        test_extraction()
    else:
        print("\nNo documents created. Please install required libraries:")
        print("  pip install python-docx python-pptx")