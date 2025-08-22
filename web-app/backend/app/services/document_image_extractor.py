"""
Document Image Extractor Service
Extracts images from Office documents and PDFs, then applies OCR
"""
import os
import io
import tempfile
import logging
from typing import List, Dict, Any, Optional
from PIL import Image
import base64

logger = logging.getLogger(__name__)

# Import document processing libraries
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available")

try:
    import openpyxl
    from openpyxl.drawing.image import Image as XlsxImage
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False
    logger.warning("openpyxl not available")

try:
    import PyPDF2
    from pdf2image import convert_from_path
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logger.warning("PyPDF2/pdf2image not available")


class DocumentImageExtractor:
    """Extract images from various document formats"""
    
    def __init__(self, ocr_service=None):
        """
        Initialize with optional OCR service
        
        Args:
            ocr_service: OCR service instance (PaddleOCR or Mock)
        """
        self.ocr_service = ocr_service
        self.temp_dir = tempfile.gettempdir()
    
    def extract_from_docx(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract images from Word document
        
        Returns:
            List of image data with OCR results
        """
        if not DOCX_AVAILABLE:
            logger.warning("python-docx not available for Word document processing")
            return []
        
        extracted_images = []
        
        try:
            doc = DocxDocument(file_path)
            
            # Extract images from document relationships
            for i, rel in enumerate(doc.part.rels.values()):
                if "image" in rel.reltype:
                    try:
                        image_data = rel.target_part.blob
                        image = Image.open(io.BytesIO(image_data))
                        
                        # Save temporarily for OCR
                        temp_path = os.path.join(self.temp_dir, f"docx_image_{i}.png")
                        image.save(temp_path)
                        
                        # Apply OCR if available
                        ocr_text = self._apply_ocr(temp_path)
                        
                        # Convert to base64 for preview
                        buffered = io.BytesIO()
                        image.save(buffered, format="PNG")
                        img_base64 = base64.b64encode(buffered.getvalue()).decode('utf-8')
                        
                        extracted_images.append({
                            'index': i + 1,
                            'format': image.format or 'Unknown',
                            'size': image.size,
                            'mode': image.mode,
                            'ocr_text': ocr_text,
                            'preview': f"data:image/png;base64,{img_base64[:100]}...",
                            'temp_path': temp_path  # Keep for AI analysis
                        })
                        
                        # Don't clean up yet - keep for AI analysis
                        # Cleanup will happen after all processing
                            
                    except Exception as e:
                        logger.error(f"Error processing DOCX image {i}: {e}")
                        
        except Exception as e:
            logger.error(f"Error extracting from DOCX: {e}")
        
        return extracted_images
    
    def extract_from_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract images from PowerPoint presentation
        
        Returns:
            List of image data with OCR results
        """
        if not PPTX_AVAILABLE:
            logger.warning("python-pptx not available for PowerPoint processing")
            return []
        
        extracted_images = []
        
        try:
            prs = Presentation(file_path)
            image_count = 0
            
            for slide_num, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        try:
                            image_data = shape.image.blob
                            image = Image.open(io.BytesIO(image_data))
                            
                            # Save temporarily for OCR
                            temp_path = os.path.join(self.temp_dir, f"pptx_image_{image_count}.png")
                            image.save(temp_path)
                            
                            # Apply OCR if available
                            ocr_text = self._apply_ocr(temp_path)
                            
                            # Convert to base64 for preview
                            buffered = io.BytesIO()
                            image.save(buffered, format="PNG")
                            img_base64 = base64.b64encode(buffered.getvalue()).decode('utf-8')
                            
                            extracted_images.append({
                                'slide': slide_num,
                                'index': image_count + 1,
                                'format': image.format or 'Unknown',
                                'size': image.size,
                                'mode': image.mode,
                                'ocr_text': ocr_text,
                                'preview': f"data:image/png;base64,{img_base64[:100]}...",
                                'temp_path': temp_path  # Keep for AI analysis
                            })
                            
                            image_count += 1
                            
                            # Keep temp file for AI analysis - will be cleaned up later
                                
                        except Exception as e:
                            logger.error(f"Error processing PPTX image on slide {slide_num}: {e}")
                            
        except Exception as e:
            logger.error(f"Error extracting from PPTX: {e}")
        
        return extracted_images
    
    def extract_from_xlsx(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract images from Excel spreadsheet
        
        Returns:
            List of image data with OCR results
        """
        if not XLSX_AVAILABLE:
            logger.warning("openpyxl not available for Excel processing")
            return []
        
        extracted_images = []
        
        try:
            wb = openpyxl.load_workbook(file_path)
            image_count = 0
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                
                # Check for images in the sheet
                if hasattr(sheet, '_images'):
                    for img in sheet._images:
                        try:
                            # Extract image data
                            image_data = img._data()
                            image = Image.open(io.BytesIO(image_data))
                            
                            # Save temporarily for OCR
                            temp_path = os.path.join(self.temp_dir, f"xlsx_image_{image_count}.png")
                            image.save(temp_path)
                            
                            # Apply OCR if available
                            ocr_text = self._apply_ocr(temp_path)
                            
                            # Convert to base64 for preview
                            buffered = io.BytesIO()
                            image.save(buffered, format="PNG")
                            img_base64 = base64.b64encode(buffered.getvalue()).decode('utf-8')
                            
                            extracted_images.append({
                                'sheet': sheet_name,
                                'index': image_count + 1,
                                'format': image.format or 'Unknown',
                                'size': image.size,
                                'mode': image.mode,
                                'ocr_text': ocr_text,
                                'preview': f"data:image/png;base64,{img_base64[:100]}...",
                                'temp_path': temp_path  # Keep for AI analysis
                            })
                            
                            image_count += 1
                            
                            # Keep temp file for AI analysis - will be cleaned up later
                                
                        except Exception as e:
                            logger.error(f"Error processing XLSX image in sheet {sheet_name}: {e}")
                            
        except Exception as e:
            logger.error(f"Error extracting from XLSX: {e}")
        
        return extracted_images
    
    def extract_from_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract images from PDF document
        
        Returns:
            List of image data with OCR results
        """
        if not PDF_AVAILABLE:
            logger.warning("PyPDF2/pdf2image not available for PDF processing")
            return []
        
        extracted_images = []
        
        try:
            # Convert PDF pages to images for OCR
            # This requires poppler-utils to be installed
            try:
                images = convert_from_path(file_path, dpi=200)
                
                for i, image in enumerate(images):
                    # Save temporarily for OCR
                    temp_path = os.path.join(self.temp_dir, f"pdf_page_{i}.png")
                    image.save(temp_path)
                    
                    # Apply OCR if available
                    ocr_text = self._apply_ocr(temp_path)
                    
                    # Create thumbnail for preview
                    thumbnail = image.copy()
                    thumbnail.thumbnail((200, 200))
                    buffered = io.BytesIO()
                    thumbnail.save(buffered, format="PNG")
                    img_base64 = base64.b64encode(buffered.getvalue()).decode('utf-8')
                    
                    extracted_images.append({
                        'page': i + 1,
                        'type': 'page_as_image',
                        'size': image.size,
                        'ocr_text': ocr_text,
                        'preview': f"data:image/png;base64,{img_base64[:100]}...",
                        'temp_path': temp_path  # Keep for AI analysis
                    })
                    
                    # Keep temp file for AI analysis - will be cleaned up later
                        
            except Exception as e:
                logger.warning(f"Could not convert PDF pages to images: {e}")
                logger.info("Install poppler-utils for PDF page conversion: sudo apt-get install poppler-utils")
                
        except Exception as e:
            logger.error(f"Error extracting from PDF: {e}")
        
        return extracted_images
    
    def _apply_ocr(self, image_path: str) -> str:
        """
        Apply OCR to an image file
        
        Args:
            image_path: Path to the image file
            
        Returns:
            Extracted text or empty string
        """
        if not self.ocr_service:
            return ""
        
        try:
            # Use the OCR service to extract text
            if hasattr(self.ocr_service, 'extract_text'):
                text = self.ocr_service.extract_text(image_path)
                return text if text else ""
            else:
                return ""
        except Exception as e:
            logger.error(f"OCR error on {image_path}: {e}")
            return ""
    
    def extract_all_images(self, file_path: str) -> Dict[str, Any]:
        """
        Extract images from any supported document type
        
        Args:
            file_path: Path to the document
            
        Returns:
            Dictionary with extracted images and metadata
        """
        file_ext = os.path.splitext(file_path)[1].lower()
        
        result = {
            'file': os.path.basename(file_path),
            'type': file_ext[1:] if file_ext else 'unknown',
            'images': [],
            'total_images': 0,
            'has_text': False
        }
        
        if file_ext == '.docx':
            result['images'] = self.extract_from_docx(file_path)
        elif file_ext == '.pptx':
            result['images'] = self.extract_from_pptx(file_path)
        elif file_ext == '.xlsx':
            result['images'] = self.extract_from_xlsx(file_path)
        elif file_ext == '.pdf':
            result['images'] = self.extract_from_pdf(file_path)
        else:
            logger.warning(f"Unsupported file type: {file_ext}")
        
        result['total_images'] = len(result['images'])
        result['has_text'] = any(
            img.get('ocr_text', '').strip() 
            for img in result['images']
        )
        
        return result
    
    def cleanup_temp_files(self, image_data_list: List[Dict[str, Any]]):
        """
        Clean up temporary files after processing
        
        Args:
            image_data_list: List of image data with temp_path keys
        """
        for img_data in image_data_list:
            temp_path = img_data.get('temp_path')
            if temp_path and os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                    logger.debug(f"Cleaned up temp file: {temp_path}")
                except Exception as e:
                    logger.warning(f"Could not remove temp file {temp_path}: {e}")