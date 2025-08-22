#!/usr/bin/env python3
"""Create a test PPTX file"""

from pptx import Presentation

# Create a new presentation
prs = Presentation()

# Add title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Test Presentation"
subtitle.text = "This is a test PowerPoint file"

# Add content slide
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Main Content"
tf = body_shape.text_frame
tf.text = "First bullet point"
p = tf.add_paragraph()
p.text = "Second bullet point"
p.level = 0
p = tf.add_paragraph()
p.text = "Third bullet point"
p.level = 0

# Save the presentation
prs.save("test_presentation.pptx")
print("Created test_presentation.pptx")