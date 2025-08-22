#!/usr/bin/env python3
"""Test conversion for all supported formats"""

import asyncio
import httpx
import tempfile
import os
from docx import Document
from pptx import Presentation
from openpyxl import Workbook
from fpdf import FPDF

async def test_format(file_path, file_name, content_type=None):
    """Test a specific file format"""
    async with httpx.AsyncClient() as client:
        with open(file_path, 'rb') as f:
            files = {"file": (file_name, f, content_type or "application/octet-stream")}
            response = await client.post(
                "http://localhost:8000/api/v1/conversion/upload",
                files=files
            )
        
        if response.status_code == 200:
            result = response.json()
            print(f"✓ {file_name}: {result.get('status')}")
            if result.get('markdown_content'):
                preview = result['markdown_content'][:100].replace('\n', ' ')
                print(f"  Preview: {preview}...")
            return True
        else:
            print(f"✗ {file_name}: Error {response.status_code}")
            print(f"  {response.text}")
            return False

async def main():
    """Test all formats"""
    print("Testing file format conversions\n")
    
    # Create test DOCX
    doc = Document()
    doc.add_heading('Test Document', 0)
    doc.add_paragraph('This is a test Word document.')
    doc.add_heading('Section 1', level=1)
    doc.add_paragraph('Some content here.')
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        docx_path = tmp.name
        doc.save(docx_path)
    
    # Create test PPTX
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Test Presentation"
    title_slide.placeholders[1].text = "Subtitle"
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        pptx_path = tmp.name
        prs.save(pptx_path)
    
    # Create test XLSX
    wb = Workbook()
    ws = wb.active
    ws.title = "Test Sheet"
    ws['A1'] = 'Name'
    ws['B1'] = 'Value'
    ws['A2'] = 'Test'
    ws['B2'] = 123
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        xlsx_path = tmp.name
        wb.save(xlsx_path)
    
    # Create test TXT
    with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False) as tmp:
        txt_path = tmp.name
        tmp.write("# Test Text File\n\nThis is plain text content.\n")
    
    # Test conversions
    results = []
    
    print("1. Testing DOCX conversion:")
    results.append(await test_format(docx_path, "test.docx", 
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"))
    
    print("\n2. Testing PPTX conversion:")
    results.append(await test_format(pptx_path, "test.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"))
    
    print("\n3. Testing XLSX conversion:")
    results.append(await test_format(xlsx_path, "test.xlsx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"))
    
    print("\n4. Testing TXT conversion:")
    results.append(await test_format(txt_path, "test.txt", "text/plain"))
    
    # Clean up
    for path in [docx_path, pptx_path, xlsx_path, txt_path]:
        try:
            os.unlink(path)
        except:
            pass
    
    # Summary
    print("\n" + "="*50)
    print(f"Results: {sum(results)}/{len(results)} formats converted successfully")
    
    return all(results)

if __name__ == "__main__":
    success = asyncio.run(main())
    exit(0 if success else 1)