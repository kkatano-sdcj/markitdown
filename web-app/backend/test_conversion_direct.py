#!/usr/bin/env python3
"""Test conversion service directly"""

import sys
import os
import asyncio
import logging

# Add app directory to path
sys.path.insert(0, '/home/kkatano/projects/markitdown/web-app/backend')

# Set up logging
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

async def test_conversion():
    """Test the conversion service directly"""
    try:
        # Import after path setup
        from app.services.conversion_service import ConversionService
        from pptx import Presentation
        import tempfile
        
        # Create a test PPTX file
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Test Presentation"
        subtitle.text = "This is a test subtitle"
        
        # Save to temp file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
            pptx_path = tmp_file.name
            prs.save(pptx_path)
            print(f"Created test file: {pptx_path}")
        
        # Create output directories
        os.makedirs("./uploads", exist_ok=True)
        os.makedirs("./converted", exist_ok=True)
        
        # Test conversion
        service = ConversionService()
        result = await service.convert_file(pptx_path, "test_output.md", save_to_db=False)
        
        print(f"\nConversion Result:")
        print(f"  Status: {result.status}")
        print(f"  Input: {result.input_file}")
        print(f"  Output: {result.output_file}")
        if result.error_message:
            print(f"  Error: {result.error_message}")
        if result.markdown_content:
            print(f"\nMarkdown Content (first 500 chars):")
            print("-" * 40)
            print(result.markdown_content[:500])
            print("-" * 40)
        
        # Clean up
        os.unlink(pptx_path)
        
        return result.status.value == "completed"
        
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    success = asyncio.run(test_conversion())
    exit(0 if success else 1)